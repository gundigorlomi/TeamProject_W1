"""Generate the Veracity project presentation (.pptx) with slide transitions
and per-shape fade-in entrance animations.

Run from the repo root:
    python generate_deck.py
Outputs: Veracity_Presentation.pptx
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn, nsmap
from lxml import etree


# ── palette ────────────────────────────────────────────────────────────────
INK = RGBColor(0x0F, 0x17, 0x2A)         # slate-900
INK_SOFT = RGBColor(0x33, 0x41, 0x55)    # slate-700
MUTED = RGBColor(0x64, 0x74, 0x8B)       # slate-500
LINE = RGBColor(0xE2, 0xE8, 0xF0)        # slate-200
WASH = RGBColor(0xF8, 0xFA, 0xFC)        # slate-50
CARD = RGBColor(0xFF, 0xFF, 0xFF)
INDIGO = RGBColor(0x63, 0x66, 0xF1)
INDIGO_DEEP = RGBColor(0x43, 0x38, 0xCA)
PINK = RGBColor(0xEC, 0x48, 0x99)
EMERALD = RGBColor(0x10, 0xB9, 0x81)
AMBER = RGBColor(0xF5, 0x9E, 0x0B)
ROSE = RGBColor(0xEF, 0x44, 0x44)
ACCENT_BG = RGBColor(0xEE, 0xEF, 0xFE)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# ── animation helpers ──────────────────────────────────────────────────────
def add_fade_transition(slide, speed: str = "med"):
    """Append <p:transition spd="med"><p:fade/></p:transition> to slide XML."""
    sld = slide._element
    # remove any existing transition
    for old in sld.findall(qn("p:transition")):
        sld.remove(old)
    transition = etree.SubElement(sld, qn("p:transition"))
    transition.set("spd", speed)
    etree.SubElement(transition, qn("p:fade"))


def add_entrance_animations(slide, shape_ids: list[int], stagger_ms: int = 250, dur_ms: int = 600):
    """Build a <p:timing> block that fades in each shape sequentially after the
    slide loads. Stagger is between shapes; duration is the fade time per shape.

    Note: removes any existing <p:timing> for the slide.
    """
    if not shape_ids:
        return
    sld = slide._element
    for old in sld.findall(qn("p:timing")):
        sld.remove(old)

    # Build the XML as a string for simplicity, then parse it.
    # Each <p:par> at the cTn(id=4) level is one "click" group; we put all
    # shapes into one click group with sequential afterEffect timing.
    next_id = [10]

    def nid() -> int:
        v = next_id[0]
        next_id[0] += 1
        return v

    inner_pars = []
    delay_acc = 0
    for sid in shape_ids:
        c5 = nid()
        c6 = nid()
        c7 = nid()
        inner_pars.append(
            f"""
            <p:par>
              <p:cTn id="{c5}" presetID="10" presetClass="entr" presetSubtype="0"
                     fill="hold" grpId="0" nodeType="afterEffect">
                <p:stCondLst><p:cond delay="{delay_acc}"/></p:stCondLst>
                <p:childTnLst>
                  <p:set>
                    <p:cBhvr>
                      <p:cTn id="{c6}" dur="1" fill="hold">
                        <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                      </p:cTn>
                      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:to><p:strVal val="visible"/></p:to>
                  </p:set>
                  <p:anim calcmode="lin" valueType="num">
                    <p:cBhvr>
                      <p:cTn id="{c7}" dur="{dur_ms}" fill="hold"/>
                      <p:tgtEl><p:spTgt spid="{sid}"/></p:tgtEl>
                      <p:attrNameLst><p:attrName>style.opacity</p:attrName></p:attrNameLst>
                    </p:cBhvr>
                    <p:tavLst>
                      <p:tav tm="0"><p:val><p:fltVal val="0"/></p:val></p:tav>
                      <p:tav tm="100000"><p:val><p:fltVal val="1"/></p:val></p:tav>
                    </p:tavLst>
                  </p:anim>
                </p:childTnLst>
              </p:cTn>
            </p:par>"""
        )
        delay_acc += stagger_ms

    inner_xml = "".join(inner_pars)

    timing_xml = f"""<p:timing xmlns:p="{P_NS}">
      <p:tnLst>
        <p:par>
          <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
            <p:childTnLst>
              <p:seq concurrent="1" nextAc="seek">
                <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
                  <p:childTnLst>
                    <p:par>
                      <p:cTn id="3" fill="hold">
                        <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
                        <p:childTnLst>
                          <p:par>
                            <p:cTn id="4" fill="hold">
                              <p:stCondLst><p:cond delay="0"/></p:stCondLst>
                              <p:childTnLst>{inner_xml}</p:childTnLst>
                            </p:cTn>
                          </p:par>
                        </p:childTnLst>
                      </p:cTn>
                    </p:par>
                  </p:childTnLst>
                </p:cTn>
                <p:prevCondLst>
                  <p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
                </p:prevCondLst>
                <p:nextCondLst>
                  <p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond>
                </p:nextCondLst>
              </p:seq>
            </p:childTnLst>
          </p:cTn>
        </p:par>
      </p:tnLst>
    </p:timing>"""

    timing_el = etree.fromstring(timing_xml)
    sld.append(timing_el)


# ── shape helpers ──────────────────────────────────────────────────────────
def fill_solid(shape, rgb: RGBColor):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb


def no_line(shape):
    shape.line.fill.background()


def set_line(shape, rgb: RGBColor, width_pt: float = 0.75):
    shape.line.color.rgb = rgb
    shape.line.width = Pt(width_pt)


def add_text(shape, text: str, *, size: int, color: RGBColor, bold: bool = False,
             align=PP_ALIGN.LEFT, font: str = "Inter", anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.text = text
    r = p.runs[0]
    r.font.name = font
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color


def add_textbox(slide, x, y, w, h, text, *, size, color, bold=False,
                align=PP_ALIGN.LEFT, font="Inter", anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(x, y, w, h)
    add_text(box, text, size=size, color=color, bold=bold, align=align, font=font, anchor=anchor)
    return box


def add_rect(slide, x, y, w, h, *, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_w: float = 0.75,
             rounded: bool = False, corner: float = 0.05):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE, x, y, w, h
    )
    if rounded:
        # adjust the corner radius (0.0 - 0.5)
        try:
            shp.adjustments[0] = corner
        except Exception:
            pass
    if fill is not None:
        fill_solid(shp, fill)
    else:
        shp.fill.background()
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_w)
    return shp


def add_oval(slide, x, y, w, h, *, fill: RGBColor | None = None,
             line: RGBColor | None = None, line_w: float = 0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, w, h)
    if fill is not None:
        fill_solid(shp, fill)
    else:
        shp.fill.background()
    if line is None:
        no_line(shp)
    else:
        set_line(shp, line, line_w)
    return shp


def chip(slide, x, y, text, *, fill=ACCENT_BG, fg=INDIGO_DEEP, size=11):
    """Small pill with text."""
    # measure-ish width: scale by char count for monospace-ish text
    w = Inches(0.18 + 0.085 * max(2, len(text)))
    h = Inches(0.32)
    shp = add_rect(slide, x, y, w, h, fill=fill, rounded=True, corner=0.5)
    tf = shp.text_frame
    tf.margin_left = Emu(60000)
    tf.margin_right = Emu(60000)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    p.text = text
    r = p.runs[0]
    r.font.name = "Inter"
    r.font.size = Pt(size)
    r.font.bold = True
    r.font.color.rgb = fg
    return shp


# ── layout primitives ─────────────────────────────────────────────────────
def blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank


def add_background(slide):
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=WASH)


def add_chrome(slide, page_num: int, total: int, eyebrow: str | None = None):
    """Top accent bar + footer page number + brand. Returns nothing tracked."""
    # subtle top accent
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.08), fill=INDIGO)

    # brand mark top-left
    add_rect(slide, Inches(0.6), Inches(0.32), Inches(0.32), Inches(0.32),
             fill=INK, rounded=True, corner=0.25)
    add_textbox(slide, Inches(1.0), Inches(0.32), Inches(2.5), Inches(0.32),
                "Veracity", size=12, color=INK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(1.0), Inches(0.55), Inches(3.0), Inches(0.22),
                "Influencer Authenticity", size=8, color=MUTED, anchor=MSO_ANCHOR.TOP)

    # eyebrow chip top-right
    if eyebrow:
        chip(slide, Inches(11.0), Inches(0.4), eyebrow, fill=ACCENT_BG, fg=INDIGO_DEEP, size=9)

    # footer
    add_rect(slide, Inches(0.6), Inches(7.05), Inches(12.13), Emu(9525), fill=LINE)
    add_textbox(slide, Inches(0.6), Inches(7.13), Inches(6), Inches(0.3),
                "Veracity · Influencer Authenticity", size=9, color=MUTED)
    add_textbox(slide, Inches(11.5), Inches(7.13), Inches(1.3), Inches(0.3),
                f"{page_num:02d} / {total:02d}", size=9, color=MUTED, align=PP_ALIGN.RIGHT,
                font="Consolas")


def slide_title(slide, title: str, eyebrow: str = "", *, y=Inches(1.05)):
    if eyebrow:
        add_textbox(slide, Inches(0.6), y, Inches(8), Inches(0.3),
                    eyebrow.upper(), size=10, color=INDIGO, bold=True)
        y2 = y + Inches(0.36)
    else:
        y2 = y
    title_box = add_textbox(slide, Inches(0.6), y2, Inches(12), Inches(0.9),
                            title, size=34, color=INK, bold=True)
    # underline accent
    add_rect(slide, Inches(0.6), y2 + Inches(0.95), Inches(0.6), Emu(34000), fill=INDIGO)
    return title_box


def shape_id(shape) -> int:
    """Return the integer shape id used for animation targeting."""
    sp = shape._element
    return int(sp.find(qn("p:nvSpPr") + "/" + qn("p:cNvPr")).get("id"))


# ── slide builders ─────────────────────────────────────────────────────────
def build_cover(prs, total):
    s = blank_slide(prs)
    add_background(s)

    # left-side gradient panel via stacked rectangles (faux gradient)
    add_rect(s, 0, 0, Inches(5.4), SLIDE_H, fill=INK)
    # decorative ovals
    big = add_oval(s, Inches(-2.0), Inches(-1.5), Inches(5.5), Inches(5.5), fill=INDIGO_DEEP)
    big.fill.transparency = 0  # solid; we'll just rely on color contrast
    add_oval(s, Inches(2.0), Inches(4.5), Inches(2.4), Inches(2.4), fill=PINK)
    add_oval(s, Inches(3.0), Inches(0.6), Inches(1.4), Inches(1.4), fill=EMERALD)

    # left panel content
    chip(s, Inches(0.6), Inches(0.6), "TEAM PROJECT · W1", fill=RGBColor(0x1E, 0x29, 0x3B), fg=RGBColor(0xC7, 0xD2, 0xFE), size=10)
    title = add_textbox(s, Inches(0.6), Inches(2.4), Inches(4.6), Inches(2.4),
                        "Veracity", size=72, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    sub = add_textbox(s, Inches(0.6), Inches(3.7), Inches(4.6), Inches(0.7),
                      "Influencer Authenticity Checker", size=20,
                      color=RGBColor(0xC7, 0xD2, 0xFE))
    tag = add_textbox(s, Inches(0.6), Inches(4.5), Inches(4.6), Inches(2.0),
                      "Score any creator 0–100. Spot fake followers, bot engagement, "
                      "and audience–content mismatch — before the brief is signed.",
                      size=14, color=RGBColor(0xCB, 0xD5, 0xE1))

    # right panel: dashboard mock card
    card = add_rect(s, Inches(6.0), Inches(1.4), Inches(6.7), Inches(4.7),
                    fill=CARD, rounded=True, corner=0.04, line=LINE)
    add_textbox(s, Inches(6.3), Inches(1.6), Inches(6.0), Inches(0.3),
                "AUTHENTICITY REPORT", size=9, color=MUTED, bold=True)
    add_textbox(s, Inches(6.3), Inches(1.85), Inches(6.0), Inches(0.5),
                "Ivy Marlowe  ·  @ivy.glow.beauty", size=18, color=INK, bold=True)

    # gauge mock
    add_oval(s, Inches(6.5), Inches(2.7), Inches(2.6), Inches(2.6), fill=RGBColor(0xFE, 0xCA, 0xCA))
    add_oval(s, Inches(6.85), Inches(3.05), Inches(1.9), Inches(1.9), fill=CARD)
    add_textbox(s, Inches(6.85), Inches(3.45), Inches(1.9), Inches(0.8),
                "36", size=44, color=ROSE, bold=True, align=PP_ALIGN.CENTER, font="Consolas")
    add_textbox(s, Inches(6.85), Inches(4.15), Inches(1.9), Inches(0.4),
                "LIKELY FAKE", size=9, color=ROSE, bold=True, align=PP_ALIGN.CENTER)

    # subscores list
    subs = [("Follower Quality", 40, ROSE),
            ("Engagement Authenticity", 53, AMBER),
            ("Audience–Content Fit", 25, ROSE),
            ("Growth Pattern", 15, ROSE)]
    yy = Inches(2.8)
    for label, val, color in subs:
        add_textbox(s, Inches(9.4), yy, Inches(2.6), Inches(0.25),
                    label, size=10, color=INK, bold=True)
        add_textbox(s, Inches(12.1), yy, Inches(0.5), Inches(0.25),
                    str(val), size=10, color=color, bold=True, align=PP_ALIGN.RIGHT, font="Consolas")
        # bar bg
        add_rect(s, Inches(9.4), yy + Inches(0.27), Inches(3.2), Emu(50000), fill=LINE, rounded=True, corner=0.5)
        # bar fill
        add_rect(s, Inches(9.4), yy + Inches(0.27), Inches(3.2 * val / 100), Emu(50000), fill=color, rounded=True, corner=0.5)
        yy += Inches(0.6)

    add_textbox(s, Inches(0.6), Inches(6.6), Inches(5.0), Inches(0.4),
                "Team Project · Week 1", size=11, color=RGBColor(0x94, 0xA3, 0xB8), font="Consolas")

    # animations
    add_fade_transition(s)
    add_entrance_animations(s, [shape_id(title), shape_id(sub), shape_id(tag), shape_id(card)],
                            stagger_ms=300, dur_ms=700)


def bullet_slide(prs, page, total, *, eyebrow, title, intro, bullets,
                 right_panel=None):
    """Generic content slide: eyebrow + title + intro paragraph + two-column bullets,
    with an optional right-side panel renderer."""
    s = blank_slide(prs)
    add_background(s)
    add_chrome(s, page, total, eyebrow)
    title_shape = slide_title(s, title, eyebrow=eyebrow)

    intro_shape = None
    if intro:
        intro_shape = add_textbox(s, Inches(0.6), Inches(2.55), Inches(7.5), Inches(0.9),
                                   intro, size=14, color=INK_SOFT)

    # Bullets in cards (left column 6 wide)
    anim_targets = [shape_id(title_shape)]
    if intro_shape:
        anim_targets.append(shape_id(intro_shape))

    bullet_y = Inches(3.5)
    bullet_x = Inches(0.6)
    bullet_w = Inches(7.5)
    row_h = Inches(0.65)
    gap = Inches(0.18)

    for i, b in enumerate(bullets):
        y = bullet_y + i * (row_h + gap)
        # number badge
        badge = add_rect(s, bullet_x, y, Inches(0.55), row_h, fill=INDIGO,
                         rounded=True, corner=0.25)
        add_text(badge, f"{i+1:02d}", size=14, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                 align=PP_ALIGN.CENTER, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
        # body card
        card = add_rect(s, bullet_x + Inches(0.7), y, bullet_w - Inches(0.7), row_h,
                        fill=CARD, rounded=True, corner=0.10, line=LINE)
        # bullet text inside card
        tb = slide_textbox_inside_card(s, bullet_x + Inches(0.95), y, bullet_w - Inches(1.0), row_h, b)
        anim_targets.extend([shape_id(badge), shape_id(card), shape_id(tb)])

    # Right panel
    if right_panel is not None:
        ids = right_panel(s)
        anim_targets.extend(ids)

    add_fade_transition(s)
    add_entrance_animations(s, anim_targets, stagger_ms=140, dur_ms=420)


def slide_textbox_inside_card(slide, x, y, w, h, text):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # split into title & detail if "—" present
    if "—" in text:
        head, rest = text.split("—", 1)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        r1 = p.add_run()
        r1.text = head.strip()
        r1.font.name = "Inter"
        r1.font.size = Pt(13)
        r1.font.bold = True
        r1.font.color.rgb = INK
        r2 = p.add_run()
        r2.text = "  " + rest.strip()
        r2.font.name = "Inter"
        r2.font.size = Pt(12)
        r2.font.color.rgb = MUTED
    else:
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.text = text
        for r in p.runs:
            r.font.name = "Inter"
            r.font.size = Pt(13)
            r.font.color.rgb = INK
    return box


# ── right-panel renderers ─────────────────────────────────────────────────
def panel_stat_grid(stats):
    """Returns a function that draws a 2x2 grid of stat tiles on the right."""
    def render(s):
        ids = []
        x0 = Inches(8.6); y0 = Inches(2.6)
        tile_w = Inches(2.05); tile_h = Inches(2.0); gap = Inches(0.2)
        for idx, (kpi, label, accent) in enumerate(stats):
            r = idx // 2; c = idx % 2
            x = x0 + c * (tile_w + gap)
            y = y0 + r * (tile_h + gap)
            tile = add_rect(s, x, y, tile_w, tile_h, fill=CARD, rounded=True, corner=0.07, line=LINE)
            bar = add_rect(s, x, y, Inches(0.06), tile_h, fill=accent)
            no_line(bar)
            kpi_box = add_textbox(s, x + Inches(0.2), y + Inches(0.3), tile_w - Inches(0.3),
                                  Inches(0.9), kpi, size=28, color=INK, bold=True, font="Consolas")
            lbl_box = add_textbox(s, x + Inches(0.2), y + Inches(1.2), tile_w - Inches(0.3),
                                  Inches(0.7), label, size=10, color=MUTED)
            ids.extend([shape_id(tile), shape_id(bar), shape_id(kpi_box), shape_id(lbl_box)])
        return ids
    return render


def panel_quote(headline, sub):
    def render(s):
        ids = []
        card = add_rect(s, Inches(8.6), Inches(2.6), Inches(4.2), Inches(4.0),
                        fill=INK, rounded=True, corner=0.04)
        # accent stripe
        stripe = add_rect(s, Inches(8.6), Inches(2.6), Inches(0.08), Inches(4.0), fill=INDIGO)
        no_line(stripe)
        big = add_textbox(s, Inches(8.9), Inches(3.0), Inches(3.8), Inches(0.6),
                          "“", size=48, color=INDIGO, bold=True, font="Georgia")
        head = add_textbox(s, Inches(8.9), Inches(3.5), Inches(3.7), Inches(2.0),
                           headline, size=18, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True)
        subt = add_textbox(s, Inches(8.9), Inches(5.7), Inches(3.7), Inches(0.8),
                           sub, size=11, color=RGBColor(0xC7, 0xD2, 0xFE))
        ids.extend([shape_id(card), shape_id(stripe), shape_id(big), shape_id(head), shape_id(subt)])
        return ids
    return render


def panel_chips(items):
    """Right-side stack of chips with a heading."""
    def render(s):
        ids = []
        head = add_textbox(s, Inches(8.6), Inches(2.6), Inches(4.2), Inches(0.4),
                           "AT A GLANCE", size=10, color=INDIGO, bold=True)
        ids.append(shape_id(head))
        y = Inches(3.0)
        for label, sub in items:
            card = add_rect(s, Inches(8.6), y, Inches(4.2), Inches(0.7),
                            fill=CARD, rounded=True, corner=0.20, line=LINE)
            dot = add_oval(s, Inches(8.8), y + Inches(0.27), Inches(0.16), Inches(0.16), fill=INDIGO)
            lbl = add_textbox(s, Inches(9.1), y + Inches(0.08), Inches(2.5), Inches(0.3),
                              label, size=12, color=INK, bold=True)
            sub_box = add_textbox(s, Inches(9.1), y + Inches(0.36), Inches(3.6), Inches(0.3),
                                  sub, size=10, color=MUTED)
            ids.extend([shape_id(card), shape_id(dot), shape_id(lbl), shape_id(sub_box)])
            y += Inches(0.85)
        return ids
    return render


def panel_arch_diagram():
    """Three-tier architecture diagram (Frontend / Backend / Data) on the right."""
    def render(s):
        ids = []
        head = add_textbox(s, Inches(8.6), Inches(2.6), Inches(4.2), Inches(0.4),
                           "ARCHITECTURE", size=10, color=INDIGO, bold=True)
        ids.append(shape_id(head))

        tiers = [
            ("React / Vite UI", "Dashboard · Compare · Live charts", INDIGO),
            ("FastAPI Service", "Auth · Scans · Scoring engine", INK),
            ("Providers + DB", "SQLAlchemy · pluggable mock/real", EMERALD),
        ]
        y = Inches(3.05)
        for title, sub, color in tiers:
            card = add_rect(s, Inches(8.6), y, Inches(4.2), Inches(0.95),
                            fill=CARD, rounded=True, corner=0.10, line=LINE)
            stripe = add_rect(s, Inches(8.6), y, Inches(0.08), Inches(0.95), fill=color)
            no_line(stripe)
            t = add_textbox(s, Inches(8.85), y + Inches(0.13), Inches(3.9), Inches(0.35),
                            title, size=13, color=INK, bold=True)
            sb = add_textbox(s, Inches(8.85), y + Inches(0.48), Inches(3.9), Inches(0.4),
                             sub, size=10, color=MUTED)
            ids.extend([shape_id(card), shape_id(stripe), shape_id(t), shape_id(sb)])
            # connector arrow
            if y < Inches(5.0):
                arrow = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(10.55),
                                           y + Inches(0.95) + Inches(0.05),
                                           Inches(0.3), Inches(0.20))
                fill_solid(arrow, LINE); no_line(arrow)
                ids.append(shape_id(arrow))
            y += Inches(1.20)
        return ids
    return render


def panel_security_badges():
    def render(s):
        ids = []
        head = add_textbox(s, Inches(8.6), Inches(2.6), Inches(4.2), Inches(0.4),
                           "GUARDRAILS", size=10, color=INDIGO, bold=True)
        ids.append(shape_id(head))
        items = [
            ("JWT", "Stateless auth, short-lived tokens", EMERALD),
            ("Bcrypt", "Salted password hashing", INDIGO),
            ("RBAC", "Per-agency data isolation", AMBER),
            ("Audit", "Every scan logged + timestamped", INK),
        ]
        y = Inches(3.05)
        for code, desc, color in items:
            card = add_rect(s, Inches(8.6), y, Inches(4.2), Inches(0.85),
                            fill=CARD, rounded=True, corner=0.12, line=LINE)
            badge = add_rect(s, Inches(8.75), y + Inches(0.16), Inches(0.95), Inches(0.55),
                             fill=color, rounded=True, corner=0.25)
            add_text(badge, code, size=11, color=RGBColor(0xFF, 0xFF, 0xFF), bold=True,
                     align=PP_ALIGN.CENTER, font="Consolas", anchor=MSO_ANCHOR.MIDDLE)
            descbox = add_textbox(s, Inches(9.85), y + Inches(0.25), Inches(2.85), Inches(0.4),
                                  desc, size=11, color=INK_SOFT)
            ids.extend([shape_id(card), shape_id(badge), shape_id(descbox)])
            y += Inches(1.0)
        return ids
    return render


# ── content for each slide ────────────────────────────────────────────────
TOTAL = 15

SLIDES = [
    # 2. Project overview
    dict(
        eyebrow="Project Overview",
        title="What we built, in one screen",
        intro="A web app that scores any social-media creator on authenticity, "
              "broken into four signals and surfaced as a single 0–100 verdict.",
        bullets=[
            "Score — weighted composite of four authenticity signals",
            "Profile — followers, posts, age, average likes at a glance",
            "Timeline — 90-day · 6-month · 1-year engagement trends",
            "Red flags — ranked anomalies with severity tags",
        ],
        right_panel=panel_stat_grid([
            ("4", "scoring signals", INDIGO),
            ("90d", "engagement window", EMERALD),
            ("0–100", "authenticity score", AMBER),
            ("2 plat.", "Instagram + TikTok", PINK),
        ]),
    ),
    # 3. Project purpose
    dict(
        eyebrow="Purpose",
        title="Make influencer due-diligence a 10-second decision",
        intro="Today, vetting a creator means stitching together screenshots, "
              "spreadsheets, and gut feel. We compress that loop into one report.",
        bullets=[
            "Replace manual vetting — turn an hour of work into a glance",
            "Standardise the verdict — one number every stakeholder agrees on",
            "Surface the why — show the evidence, not just the score",
            "Stay extensible — pluggable providers for any social platform",
        ],
        right_panel=panel_quote(
            "Brands lose ~$1.3B / year to fake influencer engagement.",
            "Veracity gives that money a defensive moat — before contracts are signed.",
        ),
    ),
    # 4. The problem
    dict(
        eyebrow="Problem",
        title="Reach numbers lie. Brands pay for it.",
        intro="Anyone can buy followers. Engagement pods make likes look organic. "
              "By the time a campaign underperforms, the budget is already gone.",
        bullets=[
            "Bought followers — entire audiences fabricated overnight",
            "Bot engagement — generic comments, scripted likes, fake reach",
            "Audience drift — niche creator with audience nothing like the niche",
            "Velocity spikes — sudden growth that has no organic explanation",
        ],
        right_panel=panel_stat_grid([
            ("49%", "of accounts show fake activity", ROSE),
            ("$1.3B", "annual ad-fraud loss", ROSE),
            ("60s", "to buy 10k followers", AMBER),
            ("0", "industry-standard score", AMBER),
        ]),
    ),
    # 5. What problem it solves
    dict(
        eyebrow="Solution",
        title="One verdict, four signals, full evidence",
        intro="Veracity unifies the moving parts of authenticity into a single "
              "report — composite score on top, drill-downs underneath.",
        bullets=[
            "Follower quality — profile signals on the audience itself",
            "Engagement authenticity — like / comment patterns over time",
            "Audience–content fit — does the audience match the niche?",
            "Growth pattern — velocity anomalies and history checks",
        ],
        right_panel=panel_arch_diagram(),
    ),
    # 6. Core idea
    dict(
        eyebrow="Core Idea",
        title="Compose four pure scoring functions into one verdict",
        intro="Each signal is an independent, deterministic function. The engine "
              "weights and aggregates — no signal is a black box.",
        bullets=[
            "Pure functions — each subscore is testable in isolation",
            "Deterministic — same input, same score, every run",
            "Weighted composite — tunable mix without touching the UI",
            "Provider-agnostic — same engine, any platform",
        ],
        right_panel=panel_chips([
            ("follower_quality.py", "profile-level audience signals"),
            ("engagement.py", "comment + like pattern analysis"),
            ("audience_fit.py", "niche vs. demographic match"),
            ("growth.py", "follower-velocity anomaly detection"),
        ]),
    ),
    # 7. Target audience
    dict(
        eyebrow="Target Audience",
        title="Who pays for this, and why",
        intro="Any team that signs a contract with a creator. The closer the team "
              "is to the spend, the higher the willingness to pay.",
        bullets=[
            "Influencer agencies — vet rosters before pitching to brands",
            "DTC marketing teams — verify creators before media spend",
            "PR / comms — protect reputation by avoiding fake amplifiers",
            "Talent platforms — score creators at scale on sign-up",
        ],
        right_panel=panel_stat_grid([
            ("B2B", "primary go-to-market", INDIGO),
            ("SMB→Mid", "agencies of 5–200", EMERALD),
            ("Self-serve", "no sales call to start", AMBER),
            ("API", "embed in talent CRMs", PINK),
        ]),
    ),
    # 8. Value proposition
    dict(
        eyebrow="Value Proposition",
        title="Stop buying reach that doesn't exist",
        intro="Veracity pays for itself the first time it kills a bad deal. "
              "Cheap to run, expensive to skip.",
        bullets=[
            "Save spend — kill bad campaigns before the wire transfer",
            "Save time — replace a vetting spreadsheet with one URL",
            "Defend reputation — receipts for every creator you sign",
            "Scale due-diligence — score 1 or 1,000 at the same cost per scan",
        ],
        right_panel=panel_quote(
            "From an hour of guesswork to a 10-second answer.",
            "One report. Four signals. A score every stakeholder trusts.",
        ),
    ),
    # 9. Growth strategy
    dict(
        eyebrow="Growth Strategy",
        title="Land with a free score, expand into the workflow",
        intro="Acquire on a viral free check, monetise on team workflows, "
              "retain on the data network effect.",
        bullets=[
            "Free public scan — single-handle check, shareable result",
            "Team workspaces — bulk scanning, watchlists, exports",
            "API tier — embed scores into CRMs and talent platforms",
            "Network effect — every scan sharpens the model for everyone",
        ],
        right_panel=panel_chips([
            ("Land", "free single-handle check, shareable"),
            ("Expand", "team plan with bulk + watchlists"),
            ("Embed", "API tier into talent CRMs"),
            ("Retain", "scan history compounds value"),
        ]),
    ),
    # 10. Revenue model
    dict(
        eyebrow="Revenue Model",
        title="Recurring SaaS with a pay-as-you-go ceiling",
        intro="Predictable monthly revenue with usage-based upside. Margins "
              "scale with seat count; cost scales with scan volume.",
        bullets=[
            "Free — 5 scans / month, single user, public results",
            "Pro — $49 / mo, 200 scans, watchlists, CSV export",
            "Team — $199 / mo, 1k scans, multi-seat, audit log",
            "API — usage-based, $0.05 / scan beyond plan ceiling",
        ],
        right_panel=panel_stat_grid([
            ("$49", "Pro / month", INDIGO),
            ("$199", "Team / month", EMERALD),
            ("$0.05", "per scan, API", AMBER),
            ("85%", "target gross margin", PINK),
        ]),
    ),
    # 11. Costs & sustainability
    dict(
        eyebrow="Costs & Sustainability",
        title="Cheap to run today, cheap to scale tomorrow",
        intro="Stack chosen for low fixed cost. Variable cost lives in API calls "
              "to social platforms — bounded by per-plan ceilings.",
        bullets=[
            "Hosting — $20–80 / mo at MVP scale (1 small VM + managed Postgres)",
            "Storage — Postgres + S3 for scan artifacts; cents per GB",
            "Platform APIs — pay-per-call; rate-limited and cached aggressively",
            "Cache layer — repeat scans for the same handle reuse fresh data",
        ],
        right_panel=panel_chips([
            ("Hosting", "$20–80 / month at MVP scale"),
            ("DB", "managed Postgres, pennies per GB"),
            ("API calls", "cached aggressively per handle"),
            ("Background jobs", "queue isolates spikes from UX"),
        ]),
    ),
    # 12. Growth & business impact
    dict(
        eyebrow="Growth & Business Impact",
        title="Where the line goes if we get this right",
        intro="A defensible category — once a brand standardises on a score, "
              "switching costs grow with every scan stored in the workspace.",
        bullets=[
            "Year 1 — 200 paying teams, $250k ARR target",
            "Year 2 — API in 3+ talent CRMs, $1.2M ARR",
            "Moat — proprietary scan history compounds per-customer",
            "Exit lanes — strategic acquisition by an ad-fraud or talent platform",
        ],
        right_panel=panel_stat_grid([
            ("$250k", "Year 1 ARR target", EMERALD),
            ("$1.2M", "Year 2 ARR target", INDIGO),
            ("3×", "API integrations Y2", AMBER),
            ("<5%", "monthly churn target", PINK),
        ]),
    ),
    # 13. Technical structure
    dict(
        eyebrow="Technical Structure",
        title="Three layers, swappable parts",
        intro="React on top, FastAPI in the middle, SQLAlchemy + pluggable "
              "providers underneath. Each layer is replaceable on its own.",
        bullets=[
            "Frontend — Vite + React + TypeScript + Tailwind, Router-based pages",
            "Backend — FastAPI service with JWT auth and BackgroundTasks",
            "Scoring — four pure functions composed by a weighted engine",
            "Providers — Protocol-typed interface; mock today, real APIs tomorrow",
        ],
        right_panel=panel_arch_diagram(),
    ),
    # 14. System features
    dict(
        eyebrow="System Features",
        title="What the user actually gets",
        intro="Every screen exists to answer one question — is this creator's "
              "audience real? Each feature is one piece of that answer.",
        bullets=[
            "Live dashboard — gauge, subscores, animated 90d / 6m / 1y timeline",
            "Compare mode — two creators, side-by-side verdict",
            "Red-flags panel — ranked, severity-tagged anomalies",
            "Suspicious followers — sample of flagged accounts with reasons",
        ],
        right_panel=panel_chips([
            ("Dashboard", "score · subscores · timeline"),
            ("Compare", "side-by-side scan diffs"),
            ("Red flags", "severity-ranked anomalies"),
            ("Audit", "every scan timestamped"),
        ]),
    ),
    # 15. Security & ethics
    dict(
        eyebrow="Security & Ethics",
        title="Trust the score because you can audit it",
        intro="A scoring product is only as good as its accountability. "
              "We hash credentials, log every scan, and never store private data.",
        bullets=[
            "Credentials — bcrypt-hashed, never stored in plaintext",
            "Tokens — short-lived JWTs, no server-side session sprawl",
            "Public data only — no private DMs, no scraped private accounts",
            "Right to recourse — every flag is explainable; nothing is a black box",
        ],
        right_panel=panel_security_badges(),
    ),
]


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    build_cover(prs, TOTAL)

    # 2–15
    for i, cfg in enumerate(SLIDES):
        page = i + 2  # cover is page 1
        bullet_slide(prs, page, TOTAL, **cfg)

    out = "Veracity_Presentation.pptx"
    prs.save(out)
    print(f"Wrote {out} with {len(prs.slides)} slides.")


if __name__ == "__main__":
    main()
